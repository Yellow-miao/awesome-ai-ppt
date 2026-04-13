#!/usr/bin/env python3
"""
DeepPresenter 工作流 - 分层 PPTX 组装
将分层截图组装为可编辑的 PPTX
"""
import os
import re
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor


def extract_slide_title(html_path: str) -> str:
    """从 HTML 提取标题"""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            content = f.read()
        match = re.search(r'class="title"[^>]*>([^<]+)<', content)
        return match.group(1).strip() if match else Path(html_path).stem
    except:
        return Path(html_path).stem


def assemble_layered_pptx(
    screenshots_dir: str,
    slides_dir: str,
    output_path: str,
    slide_width: float = 13.333,
    slide_height: float = 7.5
):
    """
    组装分层 PPTX
    screenshots_dir: 分层截图目录
    slides_dir: HTML 文件目录（用于提取标题）
    output_path: 输出路径
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    blank_layout = prs.slide_layouts[6]  # 空白布局
    
    # 按顺序获取截图文件
    screenshot_files = sorted([
        f for f in os.listdir(screenshots_dir)
        if f.endswith('_full.png')
    ])
    
    print(f'找到 {len(screenshot_files)} 页截图')
    
    for i, screenshot in enumerate(screenshot_files):
        slide = prs.slides.add_slide(blank_layout)
        base_name = screenshot.replace('_full.png', '')
        
        # 获取对应 HTML 的标题
        html_file = os.path.join(slides_dir, f'{base_name}.html')
        title = extract_slide_title(html_file) if os.path.exists(html_file) else f'第{i+1}页'
        
        # 添加背景图
        bg_path = os.path.join(screenshots_dir, f'{base_name}_full.png')
        if os.path.exists(bg_path):
            slide.shapes.add_picture(
                bg_path, 0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # 添加文字层（透明文本框，用户可编辑）
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        
        print(f'  第 {i+1} 页: {title}')
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f'分层 PPTX 已保存: {output_path}')


def main():
    parser = argparse.ArgumentParser(description='组装分层 PPTX')
    parser.add_argument('--screenshots', '-s', required=True, help='截图目录')
    parser.add_argument('--slides', '-d', required=True, help='HTML 文件目录')
    parser.add_argument('--output', '-o', default='output/layered.pptx', help='输出文件')
    parser.add_argument('--width', type=float, default=13.333, help='幻灯片宽度（英寸）')
    parser.add_argument('--height', type=float, default=7.5, help='幻灯片高度（英寸）')
    args = parser.parse_args()
    
    assemble_layered_pptx(
        args.screenshots,
        args.slides,
        args.output,
        args.width,
        args.height
    )


if __name__ == '__main__':
    exit(main())
